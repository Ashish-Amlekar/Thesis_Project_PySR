import os
import shutil
import pypandoc
import pdfplumber
from pptx import Presentation

def ensure_pandoc():
    """Ensure pandoc is installed and available for pypandoc."""
    try:
        pypandoc.get_pandoc_path()
    except OSError:
        print("[INFO] Pandoc not found, downloading...")
        pypandoc.download_pandoc()
        print("[INFO] Pandoc downloaded successfully.")

def convert_pdf_to_md(input_path, output_path):
    """Extract text from PDF using pdfplumber."""
    try:
        with pdfplumber.open(input_path) as pdf:
            text = "\n".join([page.extract_text() or "" for page in pdf.pages])
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"[INFO] PDF extracted {input_path} → {output_path}")
        return True
    except Exception as e:
        print(f"[WARN] PDF failed {input_path}: {e}")
        return False

def convert_pptx_to_md(input_path, output_path):
    """Extract text from PPTX using python-pptx."""
    try:
        prs = Presentation(input_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        text = "\n\n".join(texts)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"[INFO] PPTX extracted {input_path} → {output_path}")
        return True
    except Exception as e:
        print(f"[WARN] PPTX failed {input_path}: {e}")
        return False

def convert_txt_to_md(input_path, output_path):
    """Copy TXT content to Markdown file."""
    try:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        with open(input_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)
        print(f"[INFO] TXT → MD: {input_path} → {output_path}")
        return True
    except Exception as e:
        print(f"[WARN] TXT failed {input_path}: {e}")
        return False

def convert_to_md(input_path: str, output_path: str):
    """Convert a file to Markdown, handling PDF/PPTX separately."""
    ext = os.path.splitext(input_path)[1].lower()

    if ext == ".md":
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        shutil.copy2(input_path, output_path)
        print(f"[INFO] Copied {input_path} → {output_path}")
        return True

    if ext == ".txt":
        return convert_txt_to_md(input_path, output_path)

    if ext == ".pdf":
        return convert_pdf_to_md(input_path, output_path)

    if ext == ".pptx":
        return convert_pptx_to_md(input_path, output_path)

    # fallback to Pandoc for other formats (docx, rtf, odt, html, etc.)
    try:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        output = pypandoc.convert_file(input_path, "md", extra_args=["--standalone"])
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(output)
        print(f"[INFO] Converted {input_path} → {output_path}")
        return True
    except Exception as e:
        print(f"[WARN] Pandoc failed for {input_path}: {e}")
        return False

def build_markdown_knowledge(raw_dir: str, kb_dir: str, clear_output: bool = True):
    """Convert all files under raw_dir into .md and save them under kb_dir."""
    ensure_pandoc()

    if clear_output and os.path.exists(kb_dir):
        shutil.rmtree(kb_dir)

    converted, copied, skipped = 0, 0, 0
    for root, _, files in os.walk(raw_dir):
        for file in files:
            in_path = os.path.join(root, file)
            rel_path = os.path.relpath(in_path, raw_dir)
            out_path = os.path.join(kb_dir, os.path.splitext(rel_path)[0] + ".md")

            if convert_to_md(in_path, out_path):
                if file.lower().endswith(".md"):
                    copied += 1
                else:
                    converted += 1
            else:
                skipped += 1

    print(f"[DONE] Markdown build → {kb_dir}")
    print(f"       Converted: {converted} | Copied md: {copied} | Skipped: {skipped}")
